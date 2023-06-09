import os
import csv
import json
from PIL import Image
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
import exifread
import filetype
import requests
from bs4 import BeautifulSoup


def extract_metadata(
        file_path,
        output_format=None,
        selected_fields=None,
        recursive=False):
    # Get the file extension
    file_extension = os.path.splitext(file_path)[1].lower()

    # Extract metadata based on file type
    if file_extension == ".jpg" or file_extension == ".jpeg":
        with open(file_path, 'rb') as f:
            tags = exifread.process_file(f, details=False)

        metadata = {}
        for tag, value in tags.items():
            if selected_fields is None or tag.lower() in selected_fields:
                metadata[tag.lower()] = str(value)

    elif file_extension == ".png" or file_extension == ".gif":
        im = Image.open(file_path)
        metadata = im.info

    elif file_extension == ".pdf":
        metadata = extract_text(file_path)

    elif file_extension == ".docx":
        document = Document(file_path)
        metadata = {}
        for paragraph in document.paragraphs:
            text = paragraph.text
            if text != "":
                metadata[text] = ""

    elif file_extension == ".pptx":
        prs = Presentation(file_path)
        metadata = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text
                        if text != "":
                            metadata.append(text)

    elif file_extension == ".html" or file_extension == ".htm":
        with open(file_path, "r", encoding="utf-8") as f:
            html = f.read()
        soup = BeautifulSoup(html, "html.parser")
        metadata = {}
        metadata["title"] = soup.title.string.strip() if soup.title else ""
        description = soup.find("meta", attrs={"name": "description"})
        metadata["description"] = description["content"].strip() if description else ""
        keywords = soup.find("meta", attrs={"name": "keywords"})
        metadata["keywords"] = keywords["content"].strip() if keywords else ""


    else:
        kind = filetype.guess(file_path)
        if kind is None:
            print("Unknown file type")
            return
        else:
            metadata = {}
            metadata["file_type"] = kind.extension

    # Output metadata
    if output_format is None or output_format == "console":
        print(metadata)
    elif output_format == "csv":
        with open('metadata.csv', mode='w', newline='') as metadata_file:
            writer = csv.writer(metadata_file)
            for key, value in metadata.items():
                writer.writerow([key, value])
    elif output_format == "json":
        with open('metadata.json', mode='w') as metadata_file:
            json.dump(metadata, metadata_file)

    return metadata


def extract_metadata_from_web(url):
    # Fetch webpage HTML
    r = requests.get(url)
    html_doc = r.text

    # Parse HTML with BeautifulSoup
    soup = BeautifulSoup(html_doc, 'html.parser')

    # Extract metadata
    metadata = {}
    metadata['title'] = soup.title.string
    for meta in soup.find_all('meta'):
        if 'name' in meta.attrs and 'content' in meta.attrs:
            name = meta.attrs['name'].lower()
            content = meta.attrs['content']
            if name != '' and content != '':
                metadata[name] = content

    return metadata
